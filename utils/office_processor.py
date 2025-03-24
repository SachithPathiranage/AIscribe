import io
import os
import tempfile
import pandas as pd
from typing import Dict, Any, List, Optional
import traceback
from PIL import Image
from utils.gemini_integration import GeminiIntegration

# For Word documents
import docx
from docx.opc.exceptions import PackageNotFoundError

# For Excel files
import openpyxl
from openpyxl.drawing.image import Image as XLImage

# For PowerPoint files
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PPTXPackageNotFoundError

def process_office_file(file, options):
    """
    Process Office files (Word, Excel, PowerPoint) to extract content
    
    Args:
        file: The uploaded office file
        options: Dictionary of analysis options
    
    Returns:
        Dict containing analysis results
    """
    # Get file name and extension
    file_name = file.name
    file_extension = os.path.splitext(file_name)[1].lower()
    
    # Initialize Gemini for analysis if available
    gemini = GeminiIntegration()
    
    # Initialize metadata
    import datetime
    timestamp = datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    
    # Initialize results with common metadata
    results = {
        'file_name': file_name,
        'file_type': 'office',
        'office_type': '',
        'timestamp': timestamp,
        'user': 'SachithPathiranage'
    }
    
    # Save uploaded file to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
        temp_file.write(file.getvalue())
        temp_path = temp_file.name
    
    try:
        # Process based on file extension
        if file_extension in ['.docx', '.doc']:
            results.update(process_word_document(temp_path, options, gemini))
            results['office_type'] = 'word'
        elif file_extension in ['.xlsx', '.xls']:
            results.update(process_excel_file(temp_path, options, gemini))
            results['office_type'] = 'excel'
        elif file_extension in ['.pptx', '.ppt']:
            results.update(process_powerpoint(temp_path, options, gemini))
            results['office_type'] = 'powerpoint'
        else:
            results['error'] = f"Unsupported office file format: {file_extension}"
        
        return results
    
    except Exception as e:
        print(f"Error processing office file: {str(e)}")
        print(traceback.format_exc())
        results['error'] = f"Error processing office file: {str(e)}"
        return results
    
    finally:
        # Clean up the temporary file
        try:
            if os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass

def process_word_document(file_path, options, gemini):
    """
    Process a Word document
    
    Args:
        file_path: Path to the temporary file
        options: Dictionary of analysis options
        gemini: GeminiIntegration instance
    
    Returns:
        Dict containing analysis results
    """
    try:
        # Load document
        doc = docx.Document(file_path)
        
        # Extract basic metadata
        results = {
            'page_count': len(doc.sections),
            'paragraph_count': len(doc.paragraphs),
            'full_text': '',
            'paragraphs': [],
            'tables': [],
            'images': [],
            'headers': []
        }
        
        # Extract all paragraphs and their text
        all_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                para_data = {
                    'text': para.text,
                    'style': para.style.name
                }
                results['paragraphs'].append(para_data)
                all_text.append(para.text)
        
        # Combine all text
        results['full_text'] = '\n'.join(all_text)
        
        # Extract tables
        for i, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            
            results['tables'].append({
                'table_index': i,
                'rows': len(table.rows),
                'columns': len(table.rows[0].cells) if table.rows else 0,
                'data': table_data
            })
        
        # Extract headers
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.style.name.startswith('Heading'):
                results['headers'].append({
                    'index': i,
                    'level': int(paragraph.style.name.replace('Heading', '')) if paragraph.style.name != 'Heading' else 1,
                    'text': paragraph.text
                })
        
        # Extract images (this is more complex and depends on your requirements)
        # This is a simplified approach - images in Word docs are stored as relationships
        try:
            image_parts = []
            # Get all related parts
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_parts.append(rel.target_part)
            
            # Process each image
            for i, img_part in enumerate(image_parts):
                image_data = img_part.blob
                # Convert to PIL Image for analysis
                try:
                    img = Image.open(io.BytesIO(image_data))
                    width, height = img.size
                    
                    # Save image data
                    img_buffer = io.BytesIO()
                    img.save(img_buffer, format=img.format if img.format else 'PNG')
                    
                    image_info = {
                        'index': i,
                        'width': width,
                        'height': height,
                        'format': img.format if img.format else 'Unknown',
                        'image': img_buffer
                    }
                    
                    # Optional: Analyze image with Gemini if available
                    if gemini and gemini.is_available() and options.get('analyze_images', True):
                        image_info['description'] = gemini.analyze_image(img, "general")
                    
                    results['images'].append(image_info)
                except Exception as img_error:
                    print(f"Error processing image {i}: {str(img_error)}")
        except Exception as img_extract_error:
            print(f"Error extracting images: {str(img_extract_error)}")
            results['image_extraction_error'] = str(img_extract_error)
        
        # Generate summary using Gemini if available
        if gemini and gemini.is_available() and options.get('generate_summary', True):
            text_to_summarize = results['full_text']
            if text_to_summarize:
                results['summary'] = gemini.analyze_text(text_to_summarize, "summarize")
                
                # Extract key points
                if options.get('extract_key_points', True):
                    results['key_points'] = gemini.analyze_text(text_to_summarize, "extract_key_points")
        
        return results
    
    except PackageNotFoundError:
        raise ValueError("Invalid or corrupted Word document")
    except Exception as e:
        print(f"Error processing Word document: {str(e)}")
        print(traceback.format_exc())
        raise

def process_excel_file(file_path, options, gemini):
    """
    Process an Excel file
    
    Args:
        file_path: Path to the temporary file
        options: Dictionary of analysis options
        gemini: GeminiIntegration instance
    
    Returns:
        Dict containing analysis results
    """
    try:
        # Load workbook
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        
        # Extract basic metadata
        results = {
            'sheet_count': len(workbook.sheetnames),
            'sheet_names': workbook.sheetnames,
            'sheets': [],
            'charts': [],
            'images': [],
            'tables': [],
            'summary': ''
        }
        
        # Extract data from each sheet
        all_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            
            # Get basic sheet info
            sheet_data = {
                'name': sheet_name,
                'dimensions': sheet.dimensions,
                'max_row': sheet.max_row,
                'max_column': sheet.max_column,
                'data_preview': [],
                'contains_formulas': False,
                'column_headers': []
            }
            
            # Get column headers if they exist (assuming first row)
            if sheet.max_row > 0:
                headers = []
                for cell in sheet[1]:
                    headers.append(str(cell.value) if cell.value is not None else "")
                sheet_data['column_headers'] = headers
            
            # Get a preview of the data (first 20 rows maximum)
            max_preview_rows = min(sheet.max_row, 20)
            
            # Create DataFrame for easier processing
            data = []
            for row in sheet.iter_rows(values_only=True, max_row=max_preview_rows):
                data.append(row)
            
            # Convert to DataFrame
            if data:
                df = pd.DataFrame(data)
                if len(data) > 1:  # Use first row as headers if available
                    df.columns = df.iloc[0]
                    df = df[1:]
                sheet_data['data_preview'] = df.to_dict('records')
                
                # Add text to overall text collection for summary
                sheet_text = []
                sheet_text.append(f"Sheet: {sheet_name}")
                for i, row in df.iterrows():
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if val])
                    if row_text:
                        sheet_text.append(row_text)
                all_text.append("\n".join(sheet_text))
                
            # Check for formulas
            for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row):
                for cell in row:
                    if cell.data_type == 'f':  # Formula
                        sheet_data['contains_formulas'] = True
                        break
                if sheet_data['contains_formulas']:
                    break
            
            # Extract images from the sheet
            if sheet._images:
                for i, img in enumerate(sheet._images):
                    try:
                        img_data = {
                            'sheet': sheet_name,
                            'index': i,
                            'description': f"Image in {sheet_name}"
                        }
                        results['images'].append(img_data)
                    except Exception as img_error:
                        print(f"Error extracting image {i} from {sheet_name}: {str(img_error)}")
            
            results['sheets'].append(sheet_data)
        
        # Generate summary using Gemini if available
        if gemini and gemini.is_available() and options.get('generate_summary', True):
            text_to_summarize = "\n\n".join(all_text)
            if text_to_summarize:
                results['summary'] = gemini.analyze_text(text_to_summarize, "summarize")
                
                # Extract key points or insights from the data
                if options.get('extract_key_points', True):
                    results['insights'] = gemini.analyze_text(
                        f"This is Excel spreadsheet data:\n{text_to_summarize}\n\nExtract key insights and patterns from this data.",
                        "analyze"
                    )
        
        return results
    
    except Exception as e:
        print(f"Error processing Excel file: {str(e)}")
        print(traceback.format_exc())
        raise

def process_powerpoint(file_path, options, gemini):
    """
    Process a PowerPoint presentation
    
    Args:
        file_path: Path to the temporary file
        options: Dictionary of analysis options
        gemini: GeminiIntegration instance
    
    Returns:
        Dict containing analysis results
    """
    try:
        # Load presentation
        presentation = Presentation(file_path)
        
        # Extract basic metadata
        results = {
            'slide_count': len(presentation.slides),
            'slides': [],
            'images': [],
            'summary': '',
            'presentation_notes': ''
        }
        
        # Process each slide
        all_text = []
        for i, slide in enumerate(presentation.slides):
            slide_data = {
                'slide_number': i + 1,
                'layout': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown',
                'title': '',
                'content': [],
                'notes': '',
                'images': []
            }
            
            # Extract title
            if slide.shapes.title:
                slide_data['title'] = slide.shapes.title.text
                all_text.append(f"Slide {i+1} Title: {slide.shapes.title.text}")
            
            # Extract text from all shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:  # Skip title as we've already extracted it
                        slide_text.append(shape.text)
                        slide_data['content'].append(shape.text)
            
            all_text.append(f"Slide {i+1} Content: {' '.join(slide_text)}")
            
            # Extract notes
            if slide.has_notes_slide and slide.notes_slide:
                if slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    slide_data['notes'] = notes_text
                    results['presentation_notes'] += f"Slide {i+1}: {notes_text}\n\n"
            
            # Extract images
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        # Extract image data
                        img_data = {
                            'slide_number': i + 1,
                            'description': f"Image on slide {i+1}"
                        }
                        slide_data['images'].append(img_data)
                        results['images'].append(img_data)
                    except Exception as img_error:
                        print(f"Error extracting image from slide {i+1}: {str(img_error)}")
            
            results['slides'].append(slide_data)
        
        # Generate summary using Gemini if available
        if gemini and gemini.is_available() and options.get('generate_summary', True):
            text_to_summarize = "\n\n".join(all_text)
            if text_to_summarize:
                results['summary'] = gemini.analyze_text(text_to_summarize, "summarize")
                
                # Extract key points
                if options.get('extract_key_points', True):
                    results['key_points'] = gemini.analyze_text(text_to_summarize, "extract_key_points")
                
                # Analyze presentation structure 
                if options.get('analyze_structure', True):
                    structure_prompt = (
                        f"This is a PowerPoint presentation with {len(presentation.slides)} slides.\n"
                        f"Slide titles and content:\n{text_to_summarize}\n\n"
                        "Please analyze the presentation structure, flow, and organization. "
                        "Identify the narrative structure and how effectively the presentation flows."
                    )
                    results['structure_analysis'] = gemini.analyze_text(structure_prompt, "analyze")
        
        return results
    
    except PPTXPackageNotFoundError:
        raise ValueError("Invalid or corrupted PowerPoint file")
    except Exception as e:
        print(f"Error processing PowerPoint file: {str(e)}")
        print(traceback.format_exc())
        raise